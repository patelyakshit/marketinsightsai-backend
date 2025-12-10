"""
Slides Service

AI-powered PowerPoint presentation generation using python-pptx.
Inspired by Manus AI's slide generation capabilities.

Features:
- AI-generated content structure
- Multiple slide layouts
- Custom themes and branding
- Image integration
- Tapestry data visualization
"""

import io
import os
import uuid
import logging
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path

from app.utils.datetime_utils import utc_now
from typing import Optional, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

from app.config import get_settings

logger = logging.getLogger(__name__)
settings = get_settings()

# Paths
BACKEND_DIR = Path(__file__).parent.parent.parent
TEMPLATES_DIR = BACKEND_DIR / "templates" / "slides"
OUTPUT_DIR = Path(settings.reports_output_path)


# =============================================================================
# Data Classes
# =============================================================================

class SlideLayout(str, Enum):
    """Available slide layouts."""
    TITLE = "title"
    TITLE_CONTENT = "title_content"
    TWO_COLUMN = "two_column"
    SECTION_HEADER = "section_header"
    BULLET_POINTS = "bullet_points"
    IMAGE_WITH_CAPTION = "image_with_caption"
    COMPARISON = "comparison"
    QUOTE = "quote"
    DATA_TABLE = "data_table"
    CHART = "chart"
    BLANK = "blank"


@dataclass
class SlideContent:
    """Content for a single slide."""
    layout: SlideLayout
    title: str
    subtitle: Optional[str] = None
    body: Optional[str] = None  # Main text content
    bullet_points: list[str] = field(default_factory=list)
    image_path: Optional[str] = None
    image_caption: Optional[str] = None
    left_column: Optional[str] = None
    right_column: Optional[str] = None
    data: Optional[dict] = None  # For tables/charts
    speaker_notes: Optional[str] = None
    metadata: dict = field(default_factory=dict)


@dataclass
class PresentationConfig:
    """Configuration for presentation generation."""
    title: str
    subtitle: Optional[str] = None
    author: str = "MarketInsightsAI"
    company: str = ""
    theme: str = "default"  # "default", "dark", "professional", "modern"
    width_inches: float = 13.333  # 16:9 widescreen
    height_inches: float = 7.5
    primary_color: str = "#155E81"  # MarketInsights blue
    secondary_color: str = "#36B37E"  # Success green
    font_family: str = "Calibri"


@dataclass
class GeneratedPresentation:
    """Result of presentation generation."""
    filename: str
    filepath: str
    slide_count: int
    file_size_bytes: int
    created_at: datetime


# =============================================================================
# Theme Colors
# =============================================================================

THEMES = {
    "default": {
        "primary": "#155E81",
        "secondary": "#36B37E",
        "accent": "#FF9500",
        "text": "#333333",
        "background": "#FFFFFF",
        "muted": "#6B778C",
    },
    "dark": {
        "primary": "#4C9AFF",
        "secondary": "#57D9A3",
        "accent": "#FFAB00",
        "text": "#E8E8E8",
        "background": "#1A1A2E",
        "muted": "#A0A0A0",
    },
    "professional": {
        "primary": "#2C3E50",
        "secondary": "#3498DB",
        "accent": "#E74C3C",
        "text": "#2C3E50",
        "background": "#FFFFFF",
        "muted": "#7F8C8D",
    },
    "modern": {
        "primary": "#6366F1",
        "secondary": "#10B981",
        "accent": "#F59E0B",
        "text": "#1F2937",
        "background": "#F9FAFB",
        "muted": "#6B7280",
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


# =============================================================================
# Slide Generation Functions
# =============================================================================

def create_presentation(config: PresentationConfig) -> Presentation:
    """Create a new presentation with the given configuration."""
    prs = Presentation()

    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(config.width_inches)
    prs.slide_height = Inches(config.height_inches)

    return prs


def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: Optional[str] = None,
    config: Optional[PresentationConfig] = None,
) -> None:
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    theme = THEMES.get(config.theme if config else "default", THEMES["default"])

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(230, 230, 230)
        subtitle_para.alignment = PP_ALIGN.CENTER


def add_section_header_slide(
    prs: Presentation,
    title: str,
    subtitle: Optional[str] = None,
    config: Optional[PresentationConfig] = None,
) -> None:
    """Add a section header slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    theme = THEMES.get(config.theme if config else "default", THEMES["default"])

    # Accent bar on left
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.3), prs.slide_height
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.8),
        Inches(11), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(theme["text"])

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(4.3),
            Inches(11), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = hex_to_rgb(theme["muted"])


def add_content_slide(
    prs: Presentation,
    title: str,
    bullet_points: list[str],
    config: Optional[PresentationConfig] = None,
) -> None:
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    theme = THEMES.get(config.theme if config else "default", THEMES["default"])

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Bullet points
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.6),
        Inches(12), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()

        para.text = f"• {point}"
        para.font.size = Pt(18)
        para.font.color.rgb = hex_to_rgb(theme["text"])
        para.space_after = Pt(12)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_content: list[str],
    right_content: list[str],
    left_header: str = "",
    right_header: str = "",
    config: Optional[PresentationConfig] = None,
) -> None:
    """Add a two-column slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    theme = THEMES.get(config.theme if config else "default", THEMES["default"])

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Left column header
    if left_header:
        left_header_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(5.5), Inches(0.5)
        )
        lh_frame = left_header_box.text_frame
        lh_para = lh_frame.paragraphs[0]
        lh_para.text = left_header
        lh_para.font.size = Pt(20)
        lh_para.font.bold = True
        lh_para.font.color.rgb = hex_to_rgb(theme["primary"])

    # Left column content
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2) if left_header else Inches(1.6),
        Inches(5.5), Inches(4.8)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, point in enumerate(left_content):
        if i == 0:
            para = left_frame.paragraphs[0]
        else:
            para = left_frame.add_paragraph()
        para.text = f"• {point}"
        para.font.size = Pt(16)
        para.font.color.rgb = hex_to_rgb(theme["text"])
        para.space_after = Pt(8)

    # Right column header
    if right_header:
        right_header_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.5),
            Inches(5.5), Inches(0.5)
        )
        rh_frame = right_header_box.text_frame
        rh_para = rh_frame.paragraphs[0]
        rh_para.text = right_header
        rh_para.font.size = Pt(20)
        rh_para.font.bold = True
        rh_para.font.color.rgb = hex_to_rgb(theme["secondary"])

    # Right column content
    right_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(2.2) if right_header else Inches(1.6),
        Inches(5.5), Inches(4.8)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, point in enumerate(right_content):
        if i == 0:
            para = right_frame.paragraphs[0]
        else:
            para = right_frame.add_paragraph()
        para.text = f"• {point}"
        para.font.size = Pt(16)
        para.font.color.rgb = hex_to_rgb(theme["text"])
        para.space_after = Pt(8)


def add_quote_slide(
    prs: Presentation,
    quote: str,
    attribution: Optional[str] = None,
    config: Optional[PresentationConfig] = None,
) -> None:
    """Add a quote slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    theme = THEMES.get(config.theme if config else "default", THEMES["default"])

    # Quote mark
    quote_mark = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(2), Inches(2)
    )
    qm_frame = quote_mark.text_frame
    qm_para = qm_frame.paragraphs[0]
    qm_para.text = "\u201C"  # Opening quote
    qm_para.font.size = Pt(120)
    qm_para.font.color.rgb = hex_to_rgb(theme["primary"])

    # Quote text
    quote_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5),
        Inches(10), Inches(3)
    )
    quote_frame = quote_box.text_frame
    quote_frame.word_wrap = True
    quote_para = quote_frame.paragraphs[0]
    quote_para.text = quote
    quote_para.font.size = Pt(28)
    quote_para.font.italic = True
    quote_para.font.color.rgb = hex_to_rgb(theme["text"])

    # Attribution
    if attribution:
        attr_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(5.5),
            Inches(10), Inches(0.5)
        )
        attr_frame = attr_box.text_frame
        attr_para = attr_frame.paragraphs[0]
        attr_para.text = f"— {attribution}"
        attr_para.font.size = Pt(18)
        attr_para.font.color.rgb = hex_to_rgb(theme["muted"])
        attr_para.alignment = PP_ALIGN.RIGHT


def add_data_slide(
    prs: Presentation,
    title: str,
    data: list[dict],
    columns: list[str],
    config: Optional[PresentationConfig] = None,
) -> None:
    """Add a data table slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    theme = THEMES.get(config.theme if config else "default", THEMES["default"])

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = hex_to_rgb(theme["primary"])
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Table
    rows = len(data) + 1  # +1 for header
    cols = len(columns)
    table_width = Inches(12)
    table_height = Inches(0.5 * rows)

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.6),
        table_width, table_height
    ).table

    # Header row
    for i, col_name in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(theme["primary"])

        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for row_idx, row_data in enumerate(data):
        for col_idx, col_name in enumerate(columns):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(row_data.get(col_name, ""))

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = hex_to_rgb(theme["text"])


# =============================================================================
# High-Level Generation Functions
# =============================================================================

async def generate_presentation_from_content(
    config: PresentationConfig,
    slides: list[SlideContent],
) -> GeneratedPresentation:
    """
    Generate a presentation from structured slide content.

    Args:
        config: Presentation configuration
        slides: List of SlideContent objects

    Returns:
        GeneratedPresentation with file details
    """
    prs = create_presentation(config)

    for slide_content in slides:
        if slide_content.layout == SlideLayout.TITLE:
            add_title_slide(prs, slide_content.title, slide_content.subtitle, config)

        elif slide_content.layout == SlideLayout.SECTION_HEADER:
            add_section_header_slide(prs, slide_content.title, slide_content.subtitle, config)

        elif slide_content.layout == SlideLayout.BULLET_POINTS:
            add_content_slide(prs, slide_content.title, slide_content.bullet_points, config)

        elif slide_content.layout == SlideLayout.TWO_COLUMN:
            left = slide_content.left_column.split("\n") if slide_content.left_column else []
            right = slide_content.right_column.split("\n") if slide_content.right_column else []
            add_two_column_slide(prs, slide_content.title, left, right, config=config)

        elif slide_content.layout == SlideLayout.QUOTE:
            add_quote_slide(prs, slide_content.body or "", slide_content.subtitle, config)

        elif slide_content.layout == SlideLayout.DATA_TABLE:
            if slide_content.data:
                columns = slide_content.data.get("columns", [])
                rows = slide_content.data.get("rows", [])
                add_data_slide(prs, slide_content.title, rows, columns, config)

        elif slide_content.layout == SlideLayout.TITLE_CONTENT:
            add_content_slide(prs, slide_content.title, slide_content.bullet_points, config)

    # Save presentation
    filename = f"presentation_{uuid.uuid4().hex[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    filepath = OUTPUT_DIR / filename

    prs.save(str(filepath))

    file_size = filepath.stat().st_size

    return GeneratedPresentation(
        filename=filename,
        filepath=str(filepath),
        slide_count=len(slides),
        file_size_bytes=file_size,
        created_at=utc_now(),
    )


async def generate_tapestry_presentation(
    store_name: str,
    location: str,
    segments: list[dict],
    insights: Optional[str] = None,
    theme: str = "default",
) -> GeneratedPresentation:
    """
    Generate a Tapestry analysis presentation.

    Args:
        store_name: Name of the store/location
        location: Address or location description
        segments: List of Tapestry segment data
        insights: AI-generated insights
        theme: Presentation theme

    Returns:
        GeneratedPresentation
    """
    config = PresentationConfig(
        title=f"Tapestry Analysis: {store_name}",
        subtitle=location,
        theme=theme,
    )

    slides = []

    # Title slide
    slides.append(SlideContent(
        layout=SlideLayout.TITLE,
        title=f"Tapestry Analysis",
        subtitle=f"{store_name}\n{location}",
    ))

    # Executive summary
    if insights:
        slides.append(SlideContent(
            layout=SlideLayout.SECTION_HEADER,
            title="Executive Summary",
            subtitle="Key findings from consumer segmentation analysis",
        ))

        # Split insights into bullet points
        insight_points = [p.strip() for p in insights.split("\n") if p.strip()][:6]
        slides.append(SlideContent(
            layout=SlideLayout.BULLET_POINTS,
            title="Key Insights",
            bullet_points=insight_points,
        ))

    # Segment overview section
    slides.append(SlideContent(
        layout=SlideLayout.SECTION_HEADER,
        title="Consumer Segments",
        subtitle=f"Top {len(segments)} Tapestry segments in the trade area",
    ))

    # Segment data table
    if segments:
        segment_rows = []
        for seg in segments[:10]:  # Top 10
            segment_rows.append({
                "Segment": seg.get("name", seg.get("code", "Unknown")),
                "Code": seg.get("code", ""),
                "Percent": f"{seg.get('percent', 0):.1f}%",
                "LifeMode": seg.get("life_mode", ""),
            })

        slides.append(SlideContent(
            layout=SlideLayout.DATA_TABLE,
            title="Top Consumer Segments",
            data={
                "columns": ["Segment", "Code", "Percent", "LifeMode"],
                "rows": segment_rows,
            },
        ))

    # Individual segment slides (top 3)
    for i, seg in enumerate(segments[:3]):
        seg_name = seg.get("name", seg.get("code", f"Segment {i+1}"))
        seg_desc = seg.get("description", "")

        desc_points = [seg_desc] if seg_desc else []
        if seg.get("median_age"):
            desc_points.append(f"Median Age: {seg.get('median_age')}")
        if seg.get("median_household_income"):
            desc_points.append(f"Median HH Income: ${seg.get('median_household_income'):,}")
        if seg.get("homeownership_rate"):
            desc_points.append(f"Homeownership: {seg.get('homeownership_rate')}%")

        slides.append(SlideContent(
            layout=SlideLayout.BULLET_POINTS,
            title=f"#{i+1}: {seg_name}",
            bullet_points=desc_points or ["Detailed segment information"],
        ))

    # Closing slide
    slides.append(SlideContent(
        layout=SlideLayout.TITLE,
        title="Thank You",
        subtitle="Generated by MarketInsightsAI",
    ))

    return await generate_presentation_from_content(config, slides)


async def generate_marketing_presentation(
    campaign_name: str,
    target_audience: str,
    content_ideas: list[str],
    key_messages: list[str],
    channels: list[str],
    theme: str = "modern",
) -> GeneratedPresentation:
    """
    Generate a marketing campaign presentation.

    Args:
        campaign_name: Name of the campaign
        target_audience: Description of target audience
        content_ideas: List of content/creative ideas
        key_messages: Core messaging points
        channels: Marketing channels to use
        theme: Presentation theme

    Returns:
        GeneratedPresentation
    """
    config = PresentationConfig(
        title=campaign_name,
        subtitle="Marketing Strategy",
        theme=theme,
    )

    slides = []

    # Title
    slides.append(SlideContent(
        layout=SlideLayout.TITLE,
        title=campaign_name,
        subtitle="Marketing Strategy & Content Plan",
    ))

    # Target Audience
    slides.append(SlideContent(
        layout=SlideLayout.SECTION_HEADER,
        title="Target Audience",
        subtitle=target_audience,
    ))

    # Key Messages
    slides.append(SlideContent(
        layout=SlideLayout.BULLET_POINTS,
        title="Key Messages",
        bullet_points=key_messages,
    ))

    # Content Ideas
    slides.append(SlideContent(
        layout=SlideLayout.SECTION_HEADER,
        title="Content Strategy",
        subtitle="Creative concepts and content ideas",
    ))

    # Split content ideas across slides if many
    for i in range(0, len(content_ideas), 4):
        chunk = content_ideas[i:i+4]
        slides.append(SlideContent(
            layout=SlideLayout.BULLET_POINTS,
            title=f"Content Ideas {i//4 + 1}" if len(content_ideas) > 4 else "Content Ideas",
            bullet_points=chunk,
        ))

    # Channels
    slides.append(SlideContent(
        layout=SlideLayout.BULLET_POINTS,
        title="Marketing Channels",
        bullet_points=channels,
    ))

    # Close
    slides.append(SlideContent(
        layout=SlideLayout.TITLE,
        title="Let's Get Started",
        subtitle="Generated by MarketInsightsAI",
    ))

    return await generate_presentation_from_content(config, slides)
